"""Generate the 8-10 slide presentation deck (16:9) for the mini project.

Run from repo root:  py slides/build_slides.py
Output: slides/presentasi.pptx
"""
from __future__ import annotations
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "figures"
METRICS = json.load(open(ROOT / "metrics.json", encoding="utf-8"))

# Palette
NAVY = RGBColor(0x1F, 0x38, 0x64)
BLUE = RGBColor(0x2B, 0x6C, 0xB0)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x22, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set(tf_para, text, size, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tf_para.text = text
    tf_para.font.size = Pt(size)
    tf_para.font.bold = bold
    tf_para.font.color.rgb = color
    tf_para.alignment = align


def add_slide():
    return prs.slides.add_slide(BLANK)


def bar(slide, color=NAVY, h=1.15):
    box = slide.shapes.add_shape(1, 0, 0, SW, Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = color
    box.line.fill.background()
    box.shadow.inherit = False
    return box


def header(slide, title, kicker=None):
    bar(slide)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), SW - Inches(1.0), Inches(0.85))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf.paragraphs[0], title, 30, True, WHITE)
    if kicker:
        p = tf.add_paragraph(); _set(p, kicker, 14, False, RGBColor(0xCB, 0xDC, 0xF0))
    # page number
    pn = slide.shapes.add_textbox(SW - Inches(1.1), SH - Inches(0.55), Inches(0.9), Inches(0.4))
    _set(pn.text_frame.paragraphs[0], "", 12, False, GREY)
    return pn


def bullets(slide, items, left, top, width, height, size=18, gap=6):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, lvl, bold, color) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = "" if lvl == 0 else ("  " * lvl)
        mark = "" if bold and lvl == 0 else ("•  " if lvl == 0 else "–  ")
        _set(p, f"{prefix}{mark}{txt}", size - lvl * 2, bold, color)
        p.space_after = Pt(gap)
        p.level = lvl
    return tb


def pic(slide, name, left, top, width=None, height=None):
    kw = {}
    if width: kw["width"] = width
    if height: kw["height"] = height
    return slide.shapes.add_picture(str(FIG / name), left, top, **kw)


def caption(slide, text, left, top, width):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    _set(tb.text_frame.paragraphs[0], text, 12, False, GREY, PP_ALIGN.CENTER)


# ---------------------------------------------------------------- Slide 1 title
s = add_slide()
big = slide_bg = s.shapes.add_shape(1, 0, 0, SW, SH)
big.fill.solid(); big.fill.fore_color.rgb = NAVY; big.line.fill.background()
big.shadow.inherit = False
accent = s.shapes.add_shape(1, 0, Inches(5.0), SW, Inches(0.12))
accent.fill.solid(); accent.fill.fore_color.rgb = BLUE; accent.line.fill.background()
accent.shadow.inherit = False
tb = s.shapes.add_textbox(Inches(0.9), Inches(1.6), SW - Inches(1.8), Inches(3.2))
tf = tb.text_frame; tf.word_wrap = True
_set(tf.paragraphs[0],
     "Segmentasi Permukaan Dapat-Dipijak pada Point Cloud Anak Tangga", 36, True, WHITE)
p = tf.add_paragraph()
_set(p, "Perbandingan RANSAC, Analisis Normal PCA, DBSCAN, dan Histogram Ketinggian",
     18, False, RGBColor(0xCB, 0xDC, 0xF0))
sub = s.shapes.add_textbox(Inches(0.9), Inches(5.3), SW - Inches(1.8), Inches(1.6))
tf2 = sub.text_frame; tf2.word_wrap = True
_set(tf2.paragraphs[0], "Riyadh Lakadimu", 20, True, WHITE)
for t in ["Magister Teknik Elektro - Politeknik Elektronika Negeri Surabaya (PENS)",
          "Mata Kuliah: Sensor dan Sistem Pengolahan Sinyal (UAS)"]:
    pp = tf2.add_paragraph(); _set(pp, t, 14, False, RGBColor(0xCB, 0xDC, 0xF0))

# ---------------------------------------------------------------- Slide 2 latar
s = add_slide(); header(s, "Latar Belakang & Tujuan", "Pendahuluan")
bullets(s, [
    ("Robot pemanjat/pengikut tangga & alat bantu tunanetra butuh deteksi 3D yang andal.", 0, False, DARK),
    ("Harus bedakan permukaan dapat-dipijak (tread) dari bagian vertikal (riser) & noise.", 0, False, DARK),
    ("Salah label riser jadi 'dapat dipijak' -> pijakan ke bidang vertikal -> bahaya jatuh.", 1, False, RED),
    ("Metode deep learning butuh data berlabel besar & kurang interpretable.", 0, False, DARK),
    ("Belum ada benchmark terkontrol dengan ground-truth eksak untuk sub-masalah ini.", 1, False, DARK),
    ("Tujuan:", 0, True, NAVY),
    ("Generasi point cloud anak tangga 4-step dari model matematis berlabel.", 1, False, DARK),
    ("Implementasi & bandingkan 4 metode geometris klasik.", 1, False, DARK),
    ("Evaluasi kuantitatif + analisis noise + diskusi keamanan robot.", 1, False, DARK),
], Inches(0.6), Inches(1.45), SW - Inches(1.2), Inches(5.6), size=20)

# ---------------------------------------------------------------- Slide 3 data
s = add_slide(); header(s, "Generasi Data Point Cloud", "Model matematis 4-step")
bullets(s, [
    ("Parameter: lebar w=2.5 m, kedalaman d=0.30 m, tinggi h=0.18 m, N=4 step.", 0, False, DARK),
    ("Tread (i=0..3):", 0, True, GREEN),
    ("x~U(0,w),  y~U(i·d,(i+1)·d),  z=i·h+εz,  εz~N(0,0.02²)", 1, False, DARK),
    ("Riser (i=1..3):", 0, True, RED),
    ("x~U(0,w),  y=i·d+εy,  z~U((i-1)·h, i·h)", 1, False, DARK),
    ("+ outlier acak ~4% dan jitter LiDAR N(0,0.005²).", 0, False, DARK),
    (f"Total {METRICS['point_counts']['total']:,} titik "
     f"({METRICS['point_counts']['tread']:,} tread / "
     f"{METRICS['point_counts']['riser']:,} riser / "
     f"{METRICS['point_counts']['other']} other).", 0, True, NAVY),
], Inches(0.6), Inches(1.45), Inches(6.4), Inches(5.6), size=18)
pic(s, "profile_ground_truth.png", Inches(7.2), Inches(1.7), width=Inches(5.7))
caption(s, "Profil samping (y-z): hijau=tread, merah=riser, abu=lainnya",
        Inches(7.2), Inches(6.35), Inches(5.7))

# ---------------------------------------------------------------- Slide 4 metode overview
s = add_slide(); header(s, "Empat Metode Segmentasi", "Skema 3 kelas: tread / riser / lainnya")
bullets(s, [
    ("1. RANSAC Plane Fitting (orientation-gated)", 0, True, NAVY),
    ("Fase horizontal -> tread, fase vertikal -> riser; tolak bidang 'ramp' miring.", 1, False, DARK),
    ("2. Analisis Normal PCA", 0, True, NAVY),
    ("Normal per-titik via PCA k-NN; klasifikasi dari kemiringan |n·z|.", 1, False, DARK),
    ("3. DBSCAN + Analisis Kemiringan", 0, True, NAVY),
    ("Split orientasi -> klaster spasial -> konfirmasi slope tiap klaster.", 1, False, DARK),
    ("4. Analisis Histogram Ketinggian (lanjutan)", 0, True, NAVY),
    ("Puncak histogram z = level tread; pita antar-level = riser.", 1, False, DARK),
    ("Wajib min. 2 metode -> diimplementasikan 4 untuk perbandingan menyeluruh.", 0, True, GREEN),
], Inches(0.6), Inches(1.45), SW - Inches(1.2), Inches(5.6), size=19)

# ---------------------------------------------------------------- Slide 5 algoritma
s = add_slide(); header(s, "Metodologi & Algoritma", "Rumus inti")
bullets(s, [
    ("RANSAC: jarak titik ke bidang", 0, True, NAVY),
    ("d(x)=|n·x+b|,  inlier I={x: d(x)<τ},  τ=0.035 m", 1, False, DARK),
    ("Gate orientasi: horizontal |n·z|≥0.93 (tread), vertikal ≤0.25 (riser).", 1, False, DARK),
    ("PCA Normal: kovarians lokal k-NN (k=50)", 0, True, NAVY),
    ("C_p=(1/k)Σ(x_j-x̄)(x_j-x̄)ᵀ,  n_p=eigvec_min(C_p)", 1, False, DARK),
    ("Tread |n·z|≥0.80, riser ≤0.50 (lebih longgar: normal per-titik lebih noisy).", 1, False, DARK),
    ("DBSCAN: ε=0.08 m dalam tiap grup orientasi, lalu konfirmasi slope.", 0, True, NAVY),
    ("Bidang 'ramp' global stair: |n·z|=d/√(d²+h²)≈0.86 -> alasan gate orientasi.", 0, False, RED),
], Inches(0.6), Inches(1.45), SW - Inches(1.2), Inches(5.6), size=19)

# ---------------------------------------------------------------- Slide 6 hasil tabel
s = add_slide(); header(s, "Hasil Eksperimen", "Metrik kelas tread (one-vs-rest)")
tm = METRICS["tread_metrics"]
order = ["RANSAC", "PCA-Normals", "DBSCAN-Slope", "Height-Histogram"]
rows, cols = len(order) + 1, 6
tbl_shape = s.shapes.add_table(rows, cols, Inches(0.6), Inches(1.55),
                               Inches(7.0), Inches(2.6))
table = tbl_shape.table
hdr = ["Metode", "Acc", "Prec", "Rec", "F1", "IoU"]
for j, h in enumerate(hdr):
    c = table.cell(0, j); c.text = h
    c.fill.solid(); c.fill.fore_color.rgb = NAVY
    pr = c.text_frame.paragraphs[0]; pr.font.size = Pt(14); pr.font.bold = True
    pr.font.color.rgb = WHITE; pr.alignment = PP_ALIGN.CENTER
best = {"accuracy": "DBSCAN-Slope", "precision": "PCA-Normals", "recall": "Height-Histogram",
        "f1": "DBSCAN-Slope", "iou": "DBSCAN-Slope"}
keys = ["accuracy", "precision", "recall", "f1", "iou"]
for i, name in enumerate(order, 1):
    table.cell(i, 0).text = name
    for j, k in enumerate(keys, 1):
        v = tm[name][k]
        c = table.cell(i, j); c.text = f"{v:.3f}"
        pr = c.text_frame.paragraphs[0]; pr.font.size = Pt(13)
        pr.alignment = PP_ALIGN.CENTER
        pr.font.bold = (best[k] == name)
        if best[k] == name:
            pr.font.color.rgb = GREEN
    table.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(13)
    table.cell(i, 0).text_frame.paragraphs[0].font.bold = True
bullets(s, [
    ("DBSCAN-Slope terbaik: F1 0.943, IoU 0.892.", 0, True, GREEN),
    ("PCA-Normals dekat: F1 0.927, presisi tertinggi 0.954.", 0, False, DARK),
    ("Histogram: recall tertinggi 0.984, presisi lebih rendah.", 0, False, DARK),
    ("RANSAC murni terlemah: presisi 0.714 (riser terserap bidang tread).", 0, False, RED),
], Inches(0.6), Inches(4.4), SW - Inches(1.2), Inches(2.6), size=17)
pic(s, "seg_dbscan_slope.png", Inches(7.85), Inches(1.55), width=Inches(5.0))
caption(s, "Segmentasi DBSCAN 3D vs ground truth", Inches(7.85), Inches(3.7), Inches(5.0))

# ---------------------------------------------------------------- Slide 7 visualisasi
s = add_slide(); header(s, "Visualisasi 3D & Histogram", "Hasil eksperimen")
pic(s, "height_histogram.png", Inches(0.5), Inches(1.6), width=Inches(6.1))
caption(s, "Histogram ketinggian: 4 level tread terdeteksi (z≈0,0.18,0.36,0.54 m)",
        Inches(0.5), Inches(5.45), Inches(6.1))
pic(s, "seg_pca_normals.png", Inches(6.9), Inches(1.6), width=Inches(6.0))
caption(s, "Segmentasi PCA-Normals 3D (kiri GT, kanan hasil)",
        Inches(6.9), Inches(3.7), Inches(6.0))
bullets(s, [
    ("Tiap tread jadi level z tajam; riser mengisi pita antar-level.", 0, False, DARK),
    ("Hijau=tread, merah=riser, abu=lainnya/outlier.", 0, False, GREY),
], Inches(6.9), Inches(4.3), Inches(6.0), Inches(1.6), size=16)

# ---------------------------------------------------------------- Slide 8 noise + safety
s = add_slide(); header(s, "Analisis Noise & Keamanan Robot", "Diskusi")
pic(s, "noise_sensitivity.png", Inches(0.5), Inches(1.7), width=Inches(6.0))
caption(s, "F1/IoU tread vs σz (PCA-Normals)", Inches(0.5), Inches(5.2), Inches(6.0))
sw = METRICS["noise_sweep"]
bullets(s, [
    ("Pengaruh noise:", 0, True, NAVY),
    (f"F1 turun {sw['f1'][0]:.3f} -> {sw['f1'][-1]:.3f} saat σz 0.01 -> 0.07 m.", 1, False, DARK),
    ("Jatuh tajam saat σz mendekati ½ tinggi step (0.09 m): normal jadi ambigu.", 1, False, DARK),
    ("Keamanan robot pengikut tangga:", 0, True, NAVY),
    ("IoU≈0.89 cukup untuk perencanaan pijakan kasar.", 1, False, DARK),
    ("Presisi <1.0 -> sebagian riser salah label 'dapat dipijak' = risiko.", 1, False, RED),
    ("Mitigasi: ensemble voting, margin erosi tepi tread, verifikasi temporal.", 1, False, GREEN),
], Inches(6.8), Inches(1.5), Inches(6.1), Inches(5.6), size=18)

# ---------------------------------------------------------------- Slide 9 kesimpulan
s = add_slide(); header(s, "Kesimpulan", "Penutup")
bullets(s, [
    ("Benchmark point cloud anak tangga sintetis berlabel penuh + perbandingan 4 metode klasik.", 0, False, DARK),
    ("DBSCAN + analisis kemiringan terbaik (F1 0.943, IoU 0.892); RANSAC murni terlemah.", 0, True, NAVY),
    ("Metode berbasis normal (PCA, DBSCAN) unggul: orientasi tread vs riser cue yang kuat.", 0, False, DARK),
    ("Segmentasi memburuk saat noise mendekati ½ tinggi step.", 0, False, DARK),
    ("Cukup untuk perencanaan pijakan kasar, perlu margin keamanan sebelum dipakai robot.", 0, False, DARK),
    ("Arah ke depan: ensemble multi-metode, validasi LiDAR nyata, bandingkan vs deep learning.", 0, False, DARK),
    ("Kode + paper + repo: github.com/riyadh-gif/UAS-citraPraktikum-PakNando", 0, True, BLUE),
], Inches(0.6), Inches(1.6), SW - Inches(1.2), Inches(5.4), size=19)
tb = s.shapes.add_textbox(Inches(0.6), Inches(6.7), SW - Inches(1.2), Inches(0.6))
_set(tb.text_frame.paragraphs[0], "Terima kasih.", 20, True, NAVY)

out = Path(__file__).resolve().parent / "presentasi.pptx"
prs.save(out)
print(f"Wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
